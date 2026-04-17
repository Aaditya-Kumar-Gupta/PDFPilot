from __future__ import annotations

from pdf_toolkit.models import ToolContext, ToolError, ToolResult
from pdf_toolkit.utils.file_utils import ensure_pdf_files, resolve_output_path
from pdf_toolkit.utils.pdf_extract import extract_pages, render_page_png_stream


def run(context: ToolContext) -> ToolResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ToolError("PDF to PowerPoint requires python-pptx. Install dependencies from requirements.txt.") from exc

    ensure_pdf_files(context.input_files)
    input_file = context.input_files[0]
    output_path = resolve_output_path(context.output_path, input_file, "pdf_to_powerpoint", ".pptx")
    pages = extract_pages(input_file)
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    if pages:
        presentation.slide_width = int(pages[0].width * 12700)
        presentation.slide_height = int(pages[0].height * 12700)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    total = max(1, len(pages))
    for index, page in enumerate(pages, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(render_page_png_stream(input_file, page.page_number, dpi=200), 0, 0, width=slide_width, height=slide_height)
        context.progress(int(index / total * 100), f"Prepared PowerPoint slide for page {page.page_number}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return ToolResult(True, "PowerPoint export completed with full-slide page rendering.", [output_path], {"fallback_mode": "page-image"})
