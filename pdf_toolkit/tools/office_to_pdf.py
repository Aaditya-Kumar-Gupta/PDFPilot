from __future__ import annotations

import csv
from pathlib import Path
import subprocess

from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.utils import simpleSplit
from reportlab.pdfgen import canvas

from pdf_toolkit.models import ToolContext, ToolError, ToolResult
from pdf_toolkit.utils.renderer import resolve_office_renderer
from pdf_toolkit.utils.file_utils import ensure_files_exist, resolve_output_path
from pdf_toolkit.utils.paths import temp_dir


SOURCE_TYPE_SUFFIXES = {
    "Word": {".doc", ".docx"},
    "PowerPoint": {".ppt", ".pptx"},
    "Excel": {".xls", ".xlsx", ".xlsm", ".xltx", ".xltm", ".csv", ".tsv", ".ods", ".fods"},
}

MODERN_SUFFIXES = {
    "Word": {".docx"},
    "PowerPoint": {".pptx"},
    "Excel": {".xlsx", ".xlsm", ".xltx", ".xltm", ".csv", ".tsv"},
}

PAGE_MARGIN = 48
LINE_HEIGHT = 16


def _hidden_subprocess_kwargs() -> dict:
    if subprocess._mswindows:
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = 0
        return {
            "creationflags": subprocess.CREATE_NO_WINDOW,
            "startupinfo": startupinfo,
        }
    return {}


def _validate_source_type(input_file: Path, source_type: str, renderer_available: bool) -> None:
    expected_suffixes = SOURCE_TYPE_SUFFIXES.get(source_type)
    if not expected_suffixes:
        return
    if input_file.suffix.lower() not in expected_suffixes:
        expected = ", ".join(sorted(expected_suffixes))
        raise ToolError(f"{source_type} to PDF expects one of these file types: {expected}.")
    if renderer_available:
        return
    modern_suffixes = MODERN_SUFFIXES[source_type]
    if input_file.suffix.lower() not in modern_suffixes:
        expected = ", ".join(sorted(modern_suffixes))
        raise ToolError(f"Offline {source_type} to PDF currently supports {expected} files only.")


def _convert_with_libreoffice(soffice_path: Path, input_file: Path, output_path: Path, context: ToolContext) -> Path:
    output_dir = output_path.parent
    profile_dir = temp_dir() / "libreoffice-profile"
    profile_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.resolve().as_uri()
    command = [
        str(soffice_path),
        f"-env:UserInstallation={profile_uri}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_file),
    ]
    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        check=False,
        **_hidden_subprocess_kwargs(),
    )
    if completed.returncode != 0:
        raise ToolError(completed.stderr.strip() or "LibreOffice conversion failed.")
    produced = output_dir / f"{input_file.stem}.pdf"
    if not produced.exists():
        raise ToolError("LibreOffice did not produce the expected PDF output.")
    if produced != output_path:
        if output_path.exists():
            output_path.unlink()
        produced.replace(output_path)
    context.progress(100, "Office document converted to PDF with LibreOffice")
    return output_path


def _new_canvas(path: Path, source_type: str) -> canvas.Canvas:
    page_size = landscape(A4) if source_type in {"PowerPoint", "Excel"} else A4
    return canvas.Canvas(str(path), pagesize=page_size)


def _draw_header(pdf: canvas.Canvas, title: str, subtitle: str) -> float:
    width, height = pdf._pagesize
    y = height - PAGE_MARGIN
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(PAGE_MARGIN, y, title)
    y -= 22
    pdf.setFont("Helvetica", 10)
    pdf.drawString(PAGE_MARGIN, y, subtitle)
    return y - 20


def _write_wrapped_lines(pdf: canvas.Canvas, lines: list[str], start_y: float, font_name: str = "Helvetica", font_size: int = 11) -> None:
    width, height = pdf._pagesize
    max_width = width - PAGE_MARGIN * 2
    y = start_y
    pdf.setFont(font_name, font_size)
    for line in lines:
        parts = simpleSplit(line or " ", font_name, font_size, max_width) or [""]
        for part in parts:
            if y <= PAGE_MARGIN:
                pdf.showPage()
                pdf.setFont(font_name, font_size)
                y = height - PAGE_MARGIN
            pdf.drawString(PAGE_MARGIN, y, part)
            y -= LINE_HEIGHT


def _convert_word(input_file: Path, output_path: Path, context: ToolContext) -> None:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise ToolError("Word conversion requires python-docx.") from exc

    document = Document(input_file)
    pdf = _new_canvas(output_path, "Word")
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    if not lines:
        lines.append("This document does not contain extractable paragraph text.")
    start_y = _draw_header(pdf, input_file.name, "Converted from Word document")
    _write_wrapped_lines(pdf, lines, start_y)
    pdf.save()
    context.progress(100, "Word document converted to PDF")


def _convert_powerpoint(input_file: Path, output_path: Path, context: ToolContext) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ToolError("PowerPoint conversion requires python-pptx.") from exc

    deck = Presentation(input_file)
    pdf = _new_canvas(output_path, "PowerPoint")
    total = max(1, len(deck.slides))
    for index, slide in enumerate(deck.slides, start=1):
        title = f"{input_file.name} - Slide {index}"
        start_y = _draw_header(pdf, title, "Converted from PowerPoint presentation")
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for text_line in shape.text.splitlines():
                    clean = text_line.strip()
                    if clean:
                        lines.append(clean)
        if not lines:
            lines.append("This slide does not contain extractable text.")
        _write_wrapped_lines(pdf, lines, start_y)
        if index < total:
            pdf.showPage()
        context.progress(int(index / total * 100), f"Converted slide {index} of {total}")
    pdf.save()


def _convert_excel(input_file: Path, output_path: Path, context: ToolContext) -> None:
    pdf = _new_canvas(output_path, "Excel")
    suffix = input_file.suffix.lower()
    if suffix in {".csv", ".tsv"}:
        delimiter = "," if suffix == ".csv" else "\t"
        with input_file.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            lines = [" | ".join(str(value).strip() for value in row if str(value).strip()) for row in reader]
        lines = [line for line in lines if line]
        if not lines:
            lines = ["This file does not contain extractable row values."]
        start_y = _draw_header(pdf, input_file.name, "Converted from delimited spreadsheet")
        _write_wrapped_lines(pdf, lines, start_y, font_name="Courier", font_size=10)
        context.progress(100, "Converted delimited spreadsheet")
    else:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover
            raise ToolError("Excel conversion requires openpyxl.") from exc

        workbook = load_workbook(input_file, data_only=True)
        sheets = workbook.worksheets
        total = max(1, len(sheets))
        for index, sheet in enumerate(sheets, start=1):
            title = f"{input_file.name} - {sheet.title}"
            start_y = _draw_header(pdf, title, "Converted from Excel workbook")
            lines: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    lines.append(" | ".join(values))
            if not lines:
                lines.append("This sheet does not contain extractable cell values.")
            _write_wrapped_lines(pdf, lines, start_y, font_name="Courier", font_size=10)
            if index < total:
                pdf.showPage()
            context.progress(int(index / total * 100), f"Converted sheet {index} of {total}")
    pdf.save()


def run(context: ToolContext) -> ToolResult:
    ensure_files_exist(context.input_files)
    input_file = context.input_files[0]
    source_type = str(context.options.get("source_type", "Word") or "Word")
    renderer = resolve_office_renderer()
    _validate_source_type(input_file, source_type, renderer.available)
    output_path = resolve_output_path(context.output_path, input_file, "office_to_pdf", ".pdf")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if renderer.available and renderer.soffice_path is not None:
        rendered = _convert_with_libreoffice(renderer.soffice_path, input_file, output_path, context)
        return ToolResult(
            True,
            "Office conversion completed with high-fidelity LibreOffice rendering.",
            [rendered],
            {"source_type": source_type, "render_mode": renderer.mode},
        )

    if source_type == "Word":
        _convert_word(input_file, output_path, context)
    elif source_type == "PowerPoint":
        _convert_powerpoint(input_file, output_path, context)
    elif source_type == "Excel":
        _convert_excel(input_file, output_path, context)
    else:
        raise ToolError(f"Unsupported source type: {source_type}")

    return ToolResult(
        True,
        "Office conversion completed with built-in fallback formatting. Install or bundle LibreOffice for closer visual fidelity.",
        [output_path],
        {"source_type": source_type, "render_mode": "builtin-fallback"},
    )
