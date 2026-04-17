from __future__ import annotations

import importlib.util
import itertools
import shutil
import unittest
from pathlib import Path
from unittest.mock import patch

import fitz
from pypdf import PdfReader

from pdf_toolkit.models import ToolContext, ToolError
from pdf_toolkit.tools import office_to_pdf, pdf_to_excel, pdf_to_pdfa, pdf_to_powerpoint, pdf_to_word
from pdf_toolkit.utils.renderer import resolve_office_renderer


FIXTURES = Path(__file__).resolve().parents[1] / "runtime_tests"
TEST_OUTPUTS = Path(__file__).resolve().parents[1] / "runtime_tests" / "_test_outputs"


def _fresh_output_dir() -> Path:
    TEST_OUTPUTS.mkdir(parents=True, exist_ok=True)
    for index in itertools.count(1):
        candidate = TEST_OUTPUTS / f"run_{index}"
        if not candidate.exists():
            candidate.mkdir()
            return candidate
    raise RuntimeError("Unable to allocate test output directory.")


def _context(input_path: Path, output_path: Path) -> ToolContext:
    return ToolContext(
        input_files=[input_path],
        output_path=output_path,
        options={},
        progress=lambda *_: None,
        is_cancelled=lambda: False,
    )


def _office_context(input_path: Path, output_path: Path, source_type: str) -> ToolContext:
    return ToolContext(
        input_files=[input_path],
        output_path=output_path,
        options={"source_type": source_type},
        progress=lambda *_: None,
        is_cancelled=lambda: False,
    )


@unittest.skipUnless(importlib.util.find_spec("pdf2docx"), "pdf2docx not installed")
class PdfToWordTests(unittest.TestCase):
    def test_word_export_creates_docx(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            output = temp_dir / "sample.docx"
            result = pdf_to_word.run(_context(FIXTURES / "sample1.pdf", output))
            self.assertTrue(result.output_files[0].exists())
            self.assertEqual(result.output_files[0].suffix.lower(), ".docx")
            from docx import Document

            document = Document(result.output_files[0])
            text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
            self.assertTrue(text)
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


@unittest.skipUnless(importlib.util.find_spec("pptx"), "python-pptx not installed")
class PdfToPowerPointTests(unittest.TestCase):
    def test_powerpoint_export_creates_slide_per_page(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            output = temp_dir / "sample.pptx"
            result = pdf_to_powerpoint.run(_context(FIXTURES / "sample1.pdf", output))
            self.assertTrue(result.output_files[0].exists())
            from pptx import Presentation

            deck = Presentation(result.output_files[0])
            with fitz.open(FIXTURES / "sample1.pdf") as pdf:
                self.assertEqual(len(deck.slides), len(pdf))
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


@unittest.skipUnless(importlib.util.find_spec("openpyxl"), "openpyxl not installed")
class PdfToExcelTests(unittest.TestCase):
    def test_excel_export_creates_sheet_per_page(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            output = temp_dir / "sample.xlsx"
            result = pdf_to_excel.run(_context(FIXTURES / "sample1.pdf", output))
            self.assertTrue(result.output_files[0].exists())
            from openpyxl import load_workbook

            book = load_workbook(result.output_files[0])
            with fitz.open(FIXTURES / "sample1.pdf") as pdf:
                self.assertEqual(len(book.sheetnames), len(pdf))
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


class PdfToPdfaTests(unittest.TestCase):
    def test_archival_export_creates_unencrypted_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            output = temp_dir / "sample_pdfa.pdf"
            result = pdf_to_pdfa.run(_context(FIXTURES / "sample1.pdf", output))
            self.assertTrue(result.output_files[0].exists())
            reader = PdfReader(result.output_files[0])
            self.assertFalse(reader.is_encrypted)
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


class OfficeToPdfTests(unittest.TestCase):
    def test_word_to_pdf_creates_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            from docx import Document

            source = temp_dir / "sample.docx"
            document = Document()
            document.add_paragraph("Hello from Word")
            document.save(source)

            output = temp_dir / "word_output.pdf"
            result = office_to_pdf.run(_office_context(source, output, "Word"))
            self.assertTrue(result.output_files[0].exists())
            self.assertEqual(result.output_files[0].suffix.lower(), ".pdf")
            expected_mode = resolve_office_renderer().mode if resolve_office_renderer().available else "builtin-fallback"
            self.assertEqual(result.metadata["render_mode"], expected_mode)
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_powerpoint_to_pdf_creates_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            from pptx import Presentation

            source = temp_dir / "sample.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = "Hello slide"
            slide.placeholders[1].text = "Body text"
            deck.save(source)

            output = temp_dir / "powerpoint_output.pdf"
            result = office_to_pdf.run(_office_context(source, output, "PowerPoint"))
            self.assertTrue(result.output_files[0].exists())
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_excel_to_pdf_creates_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            from openpyxl import Workbook

            source = temp_dir / "sample.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet["A1"] = "Region"
            sheet["B1"] = "Sales"
            sheet["A2"] = "North"
            sheet["B2"] = 42
            workbook.save(source)

            output = temp_dir / "excel_output.pdf"
            result = office_to_pdf.run(_office_context(source, output, "Excel"))
            self.assertTrue(result.output_files[0].exists())
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_csv_to_pdf_creates_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            source = temp_dir / "sample.csv"
            source.write_text("Region,Sales\nNorth,42\nSouth,11\n", encoding="utf-8")

            output = temp_dir / "csv_output.pdf"
            result = office_to_pdf.run(_office_context(source, output, "Excel"))
            self.assertTrue(result.output_files[0].exists())
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_xls_is_allowed_when_renderer_is_available(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            source = temp_dir / "legacy.xls"
            source.write_bytes(b"fake-xls")
            output = temp_dir / "legacy_output.pdf"
            fake_rendered = temp_dir / "fake_legacy.pdf"
            fake_rendered.write_bytes(b"%PDF-1.4\n%%EOF\n")

            with patch("pdf_toolkit.tools.office_to_pdf.resolve_office_renderer") as mock_resolve:
                with patch("pdf_toolkit.tools.office_to_pdf._convert_with_libreoffice", return_value=fake_rendered):
                    from pdf_toolkit.utils.renderer import RendererResolution

                    mock_resolve.return_value = RendererResolution(
                        available=True,
                        mode="bundled-libreoffice",
                        soffice_path=Path(r"C:\bundle\program\soffice.exe"),
                        details="bundled",
                    )
                    result = office_to_pdf.run(_office_context(source, output, "Excel"))
            self.assertEqual(result.metadata["render_mode"], "bundled-libreoffice")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_word_to_pdf_reports_bundled_renderer_mode(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            from docx import Document

            source = temp_dir / "sample.docx"
            document = Document()
            document.add_paragraph("Hello from bundled renderer")
            document.save(source)

            output = temp_dir / "bundled_output.pdf"
            fake_rendered = temp_dir / "fake_rendered.pdf"
            fake_rendered.write_bytes(b"%PDF-1.4\n%%EOF\n")

            with patch("pdf_toolkit.tools.office_to_pdf.resolve_office_renderer") as mock_resolve:
                with patch("pdf_toolkit.tools.office_to_pdf._convert_with_libreoffice", return_value=fake_rendered):
                    from pdf_toolkit.utils.renderer import RendererResolution

                    mock_resolve.return_value = RendererResolution(
                        available=True,
                        mode="bundled-libreoffice",
                        soffice_path=Path(r"C:\bundle\program\soffice.exe"),
                        details="bundled",
                    )
                    result = office_to_pdf.run(_office_context(source, output, "Word"))
            self.assertEqual(result.metadata["render_mode"], "bundled-libreoffice")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_archival_export_rejects_encrypted_pdf(self) -> None:
        temp_dir = _fresh_output_dir()
        try:
            protected_source = temp_dir / "protected.pdf"
            shutil.copy(FIXTURES / "protected.pdf", protected_source)
            output = temp_dir / "protected_pdfa.pdf"
            with self.assertRaises(ToolError):
                pdf_to_pdfa.run(_context(protected_source, output))
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
